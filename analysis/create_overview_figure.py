#!/usr/bin/env python3
"""
Generate a large, readable overview figure + PowerPoint slide
for scorecards4extremes.
"""
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import matplotlib.patheffects as pe
import numpy as np

# ═══════════════════════════════════════════════════════════════════
#  PART 1: HIGH-RES FIGURE (large text, clean layout)
# ═══════════════════════════════════════════════════════════════════

PAL = {
    'bg':       '#F5F7FA',
    'title':    '#1A1F36',
    'sub':      '#4E5D78',
    'step':     '#1B4965',
    'step_hi':  '#5FA8D3',
    'white':    '#FFFFFF',
    'num':      '#BEE9E8',
    'det':      '#0B3954',  'det_bg': '#DBEAF5', 'det_bdr': '#6BA3CC',
    'ens':      '#7A4A1E',  'ens_bg': '#FDF0E2', 'ens_bdr': '#D4A574',
    'fc':       '#3D5A80',  'obs':    '#EE6C4D',
    'out':      '#2D6A4F',  'out_bg': '#D8F3DC',
    'cfg':      '#7B2D8E',  'cfg_bg': '#F3E8F9',
    'arr':      '#8B9DC3',
    'shadow':   '#C8D0DC',
    'bdr':      '#CDD5E0',
}


def make_figure():
    fig = plt.figure(figsize=(24, 16), facecolor=PAL['bg'])
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 24)
    ax.set_ylim(0, 16)
    ax.set_aspect('equal')
    ax.axis('off')

    # ── helpers ──────────────────────────────────────────────────
    def box(x, y, w, h, fc, ec='none', lw=0, r=0.18, z=2, shadow=False):
        if shadow:
            s = FancyBboxPatch((x+.07, y-.07), w, h,
                boxstyle=f"round,pad={r}", fc=PAL['shadow'],
                ec='none', alpha=.35, zorder=z-1)
            ax.add_patch(s)
        p = FancyBboxPatch((x, y), w, h,
            boxstyle=f"round,pad={r}", fc=fc, ec=ec, lw=lw, zorder=z)
        ax.add_patch(p)

    def step(x, y, n, label, sub='', w=4.0, h=0.95):
        box(x, y, w, h, PAL['step'], PAL['step_hi'], 2, .14, 4, True)
        c = plt.Circle((x+.52, y+h/2), .28, fc=PAL['step_hi'], ec='none', zorder=5)
        ax.add_patch(c)
        ax.text(x+.52, y+h/2, str(n), fontsize=17, fontweight='bold',
                color=PAL['step'], ha='center', va='center', zorder=6)
        dy = .08 if sub else 0
        ax.text(x+1.0, y+h/2+dy, label, fontsize=16, fontweight='bold',
                color=PAL['white'], va='center', zorder=6)
        if sub:
            ax.text(x+1.0, y+h/2-.24, sub, fontsize=11.5,
                    color=PAL['num'], va='center', zorder=6, style='italic')

    def arr_v(x, y1, y2, c=PAL['arr']):
        ax.annotate('', xy=(x, y2), xytext=(x, y1),
            arrowprops=dict(arrowstyle='->,head_width=.14,head_length=.12',
                            color=c, lw=2.8), zorder=3)

    def arr_h(x1, y, x2, c=PAL['arr']):
        ax.annotate('', xy=(x2, y), xytext=(x1, y),
            arrowprops=dict(arrowstyle='->,head_width=.14,head_length=.12',
                            color=c, lw=2.8), zorder=3)

    def arr_c(x1, y1, x2, y2, c=PAL['arr'], rad=.25):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
            arrowprops=dict(arrowstyle='->,head_width=.14,head_length=.12',
                            color=c, lw=2.8,
                            connectionstyle=f'arc3,rad={rad}'), zorder=3)

    def icard(x, y, w, h, title, lines, tc, fc, ec):
        box(x, y, w, h, fc, ec, 1.5, .14)
        ax.text(x+w/2, y+h-.35, title, fontsize=14, fontweight='bold',
                color=tc, ha='center', va='center', zorder=5)
        for i, l in enumerate(lines):
            ax.text(x+w/2, y+h-.72-i*.32, l, fontsize=11.5,
                    color=tc, ha='center', va='center', zorder=5, alpha=.75)

    def bstep(x, y, label, sub, w, fc, ec):
        box(x, y, w, .78, fc, ec, 1.8, .14, 5)
        dy = .08 if sub else 0
        ax.text(x+w/2, y+.39+dy, label, fontsize=15, fontweight='bold',
                color='#FFF', ha='center', va='center', zorder=6)
        if sub:
            ax.text(x+w/2, y+.39-.22, sub, fontsize=11,
                    color='#D0DEE8', ha='center', va='center', zorder=6, style='italic')

    # ── TITLE ────────────────────────────────────────────────────
    ax.text(12, 15.4, 'scorecards4extremes', fontsize=36, fontweight='bold',
            color=PAL['title'], ha='center', va='center', fontfamily='monospace')
    ax.text(12, 14.85, 'Configurable verification framework for extreme weather scorecard generation',
            fontsize=16, color=PAL['sub'], ha='center', va='center')
    ax.plot([2.5, 21.5], [14.55, 14.55], color='#D6DCE5', lw=1.5, zorder=1)

    # ── INPUT CARDS (left) ───────────────────────────────────────
    icard(0.4, 11.6, 4.0, 2.5, '\u25B6  Forecasts', [
        'Model 1  vs  Model 2',
        'e.g. IFS oper vs AIFS hybrid',
        'Det (control) or Ens (50 mbr)',
    ], PAL['fc'], '#E8EEF5', PAL['det_bdr'])

    icard(0.4, 8.7, 4.0, 2.5, '\u25CF  Observations', [
        'SYNOP (WMO stations)',
        'HDOBS (high-density obs)',
        'via STVL / local .gpt files',
    ], PAL['obs'], '#FFF0EC', '#E8947E')

    # ── CONFIG (top-right) ───────────────────────────────────────
    icard(19.6, 11.6, 4.0, 2.5, '\u2699  config.yaml', [
        'variable \u00B7 date range',
        'models \u00B7 thresholds \u00B7 scores',
        'filters \u00B7 bootstrap \u00B7 output',
    ], PAL['cfg'], PAL['cfg_bg'], '#C9A0DC')

    # ── BACKEND (right) ──────────────────────────────────────────
    box(19.6, 9.2, 1.85, 1.1, '#E8F5E9', '#81C784', 1.3, .12)
    ax.text(20.52, 9.93, 'Local', fontsize=13, fontweight='bold',
            color='#2D6A4F', ha='center', va='center', zorder=5)
    ax.text(20.52, 9.55, 'GRIB on disk', fontsize=10,
            color='#2D6A4F', ha='center', va='center', zorder=5, alpha=.65)

    box(21.75, 9.2, 1.85, 1.1, '#F3E5F5', '#CE93D8', 1.3, .12)
    ax.text(22.67, 9.93, 'Quaver', fontsize=13, fontweight='bold',
            color='#6A1B9A', ha='center', va='center', zorder=5)
    ax.text(22.67, 9.55, 'MARS / VTB', fontsize=10,
            color='#6A1B9A', ha='center', va='center', zorder=5, alpha=.65)

    ax.text(21.6, 9.75, '/', fontsize=18, color=PAL['sub'],
            ha='center', va='center', fontweight='bold', zorder=6)

    # ── STRATIFICATION (left, below obs) ─────────────────────────
    icard(0.4, 6.1, 4.0, 2.2, '\u25C6  Stratification', [
        'Season: DJF \u00B7 MAM \u00B7 JJA \u00B7 SON',
        'Orography: flat \u00B7 hilly \u00B7 mount.',
        'Lead time: day 1 \u2013 day 10',
    ], PAL['sub'], '#F0F2F5', PAL['bdr'])

    # ── MAIN PIPELINE ────────────────────────────────────────────
    xc = 6.5       # left edge
    sw = 4.2       # width
    gap = .38
    cx = xc + sw/2 # centre

    y = 13.2
    step(xc, y, 1, 'Read Data', 'Load forecast & observation fields', sw)
    arr_h(4.4, 12.7, xc-.15, PAL['fc'])
    arr_c(4.4, 9.8, xc-.15, y+.47, PAL['obs'], rad=-.35)
    arr_h(19.5, 12.8, xc+sw+.15, PAL['cfg'])

    y -= (.95+gap)
    step(xc, y, 2, 'Pre-process', 'Lapse-rate, wind speed, unit conversion', sw)
    arr_v(cx, 13.2, y+.95)

    y -= (.95+gap)
    step(xc, y, 3, 'Extract Points', 'Nearest gridpoint / aligned()', sw)
    arr_v(cx, y+.95+gap, y+.95)
    ax.text(xc+sw+.4, y+.47, '\u2192  Parquet files (per forecast day)',
            fontsize=12, color=PAL['out'], fontweight='bold', va='center', zorder=5)

    y -= (.95+gap)
    step(xc, y, 4, 'Filter & QC', 'Season, orography, outlier removal', sw)
    arr_v(cx, y+.95+gap, y+.95)

    y -= (.95+gap)
    step(xc, y, 5, 'Threshold', 'Fixed / dataset / station climatology', sw)
    arr_v(cx, y+.95+gap, y+.95)

    branch_y = y  # y of step 5 bottom

    # ── BRANCHES ─────────────────────────────────────────────────
    dx = 1.8; ex = 13.0
    dw = 5.2; ew = 9.2
    bh = 4.0; by_b = 1.8

    box(dx-.3, by_b, dw, bh, PAL['det_bg'], PAL['det_bdr'], 2, .22, 1, True)
    box(ex-.3, by_b, ew, bh, PAL['ens_bg'], PAL['ens_bdr'], 2, .22, 1, True)

    ax.text(dx+dw/2-.3, by_b+bh-.35, 'DETERMINISTIC', fontsize=17,
            fontweight='bold', color=PAL['det'], ha='center', va='center', zorder=5,
            path_effects=[pe.withStroke(linewidth=4, foreground=PAL['det_bg'])])
    ax.text(ex+ew/2-.3, by_b+bh-.35, 'ENSEMBLE', fontsize=17,
            fontweight='bold', color=PAL['ens'], ha='center', va='center', zorder=5,
            path_effects=[pe.withStroke(linewidth=4, foreground=PAL['ens_bg'])])

    arr_c(cx, branch_y, dx+dw/2-.3, by_b+bh-.65, PAL['det_bdr'], .3)
    arr_c(cx, branch_y, ex+ew/2-.3, by_b+bh-.65, PAL['ens_bdr'], -.3)

    # Det scores
    dcx = dx+dw/2-.3
    bstep(dx+.1, by_b+2.3, '\u2465 Scores', 'ETS \u00B7 PSS \u00B7 POD \u00B7 FAR',
          dw-.8, PAL['det'], PAL['det_bdr'])
    ax.text(dcx, by_b+2.12, 'twMAE \u00B7 twRMSE \u00B7 bias \u00B7 correlation',
            fontsize=10, color='#1A3A5C', ha='center', va='center', zorder=6, style='italic')

    bstep(dx+.1, by_b+1.05, '\u2466 Bootstrap', 'Paired significance testing',
          dw-.8, PAL['det'], PAL['det_bdr'])
    arr_v(dcx, by_b+2.3, by_b+1.05+.78, PAL['det_bdr'])

    # Ens scores
    ecx = ex+ew/2-.3
    bstep(ex+.1, by_b+2.3, '\u2465 Scores', 'twCRPS \u00B7 Brier \u00B7 BSS \u00B7 Quantile Score',
          ew-.8, '#6D4C28', PAL['ens_bdr'])
    ax.text(ecx, by_b+2.12, 'twMAE \u00B7 tw Spread-Skill Ratio \u00B7 CRPS',
            fontsize=10, color='#5C3A1A', ha='center', va='center', zorder=6, style='italic')

    bstep(ex+.1, by_b+1.05, '\u2466 Bootstrap', 'Paired significance testing',
          ew-.8, '#6D4C28', PAL['ens_bdr'])
    arr_v(ecx, by_b+2.3, by_b+1.05+.78, PAL['ens_bdr'])

    # ── BOTTOM: Save & Plot ──────────────────────────────────────
    sy = .25
    step(6.0, sy, 8, 'Save Results', 'CSV per season \u00D7 orography', 4.2)
    step(12.0, sy, 9, 'Plot Scorecards', 'Heatmaps \u00B7 panels \u00B7 diagnostics', 4.5)

    arr_c(dcx, by_b+1.05, 8.1, sy+.95+.05, PAL['det_bdr'], -.15)
    arr_c(ecx, by_b+1.05, 8.3, sy+.95+.05, PAL['ens_bdr'], .15)
    arr_h(10.2+.1, sy+.47, 12.0-.1)

    # Output
    box(17.0, sy-.05, 3.8, 1.05, PAL['out_bg'], '#52B788', 1.3, .14)
    ax.text(18.9, sy+.68, '\u2192  Scorecard Heatmaps', fontsize=13,
            fontweight='bold', color=PAL['out'], ha='center', va='center', zorder=5)
    ax.text(18.9, sy+.28, 'PNG / PDF  \u00B7  per score \u00D7 condition',
            fontsize=10.5, color=PAL['out'], ha='center', va='center', zorder=5, alpha=.7)
    arr_h(16.5+.1, sy+.47, 17.0-.1, '#52B788')

    # ── VARIABLES (bottom-left) ──────────────────────────────────
    box(0.4, .1, 4.2, 1.15, '#FFF', PAL['bdr'], 1.2, .12)
    ax.text(2.5, 1.0, 'Supported Variables', fontsize=12.5, fontweight='bold',
            color=PAL['title'], ha='center', va='center', zorder=5)
    for i, v in enumerate(['2t', '10ff', 'tp24']):
        px = .75 + i*1.35
        box(px, .25, 1.15, .48, PAL['step'], 'none', 0, .1, 5)
        ax.text(px+.575, .55, v, fontsize=12, fontweight='bold', color='#FFF',
                ha='center', va='center', zorder=6, fontfamily='monospace')

    out = './plots/scorecards4extremes_overview.png'
    fig.savefig(out, dpi=250, bbox_inches='tight', facecolor=PAL['bg'], pad_inches=.5)
    plt.close()
    print(f'Figure saved: {out}')


# ═══════════════════════════════════════════════════════════════════
#  PART 2: POWERPOINT SLIDE
# ═══════════════════════════════════════════════════════════════════

def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    def rgb(hex_str):
        h = hex_str.lstrip('#')
        return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:],16))

    def add_rounded_box(left, top, w, h, fill_hex, border_hex=None, border_w=Pt(1.5)):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left, top, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill_hex)
        if border_hex:
            shp.line.color.rgb = rgb(border_hex)
            shp.line.width = border_w
        else:
            shp.line.fill.background()
        # Smaller corner radius
        shp.adjustments[0] = 0.08
        return shp

    def add_text_box(left, top, w, h, text, font_size=14,
                     bold=False, color='#1A1F36', align=PP_ALIGN.LEFT,
                     font_name='Calibri'):
        txBox = slide.shapes.add_textbox(left, top, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
        p.font.name = font_name
        p.alignment = align
        return txBox

    def step_box(left, top, num, label, sub_text, w=Inches(3.2), h=Inches(0.7)):
        # Main box
        shp = add_rounded_box(left, top, w, h, '#1B4965', '#2B6A8E', Pt(2))
        # Number circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            left + Inches(0.12), top + Inches(0.12),
            Inches(0.46), Inches(0.46))
        circ.fill.solid()
        circ.fill.fore_color.rgb = rgb('#5FA8D3')
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].text = str(num)
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = rgb('#1B4965')
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Label
        add_text_box(left + Inches(0.65), top + Inches(0.02),
                     w - Inches(0.8), Inches(0.35),
                     label, 15, True, '#FFFFFF')
        # Sub-label
        if sub_text:
            add_text_box(left + Inches(0.65), top + Inches(0.35),
                         w - Inches(0.8), Inches(0.30),
                         sub_text, 10, False, '#BEE9E8')

    def info_box(left, top, w, h, title, lines, title_color, fill, border):
        add_rounded_box(left, top, w, h, fill, border)
        add_text_box(left + Inches(0.15), top + Inches(0.08),
                     w - Inches(0.3), Inches(0.35),
                     title, 14, True, title_color, PP_ALIGN.CENTER)
        y_off = Inches(0.42)
        for line in lines:
            add_text_box(left + Inches(0.1), top + y_off,
                         w - Inches(0.2), Inches(0.28),
                         line, 11, False, title_color, PP_ALIGN.CENTER)
            y_off += Inches(0.26)

    def add_arrow(x1, y1, x2, y2, color_hex='#8B9DC3', width=Pt(2.5)):
        """Add a connector arrow between two points."""
        # Use a freeform/line shape
        connector = slide.shapes.add_connector(
            1,  # straight connector
            x1, y1, x2, y2
        )
        connector.line.color.rgb = rgb(color_hex)
        connector.line.width = width
        # Add arrowhead
        connector.end_style = 2  # triangle arrowhead

    def branch_box(left, top, label, sub, w, fill, border):
        shp = add_rounded_box(left, top, w, Inches(0.6), fill, border, Pt(2))
        add_text_box(left + Inches(0.1), top + Inches(0.0),
                     w - Inches(0.2), Inches(0.32),
                     label, 13, True, '#FFFFFF', PP_ALIGN.CENTER)
        if sub:
            add_text_box(left + Inches(0.1), top + Inches(0.30),
                         w - Inches(0.2), Inches(0.25),
                         sub, 9.5, False, '#C8D8E8', PP_ALIGN.CENTER)

    # ── TITLE ────────────────────────────────────────────────────
    add_text_box(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 'scorecards4extremes', 32, True, '#1A1F36',
                 PP_ALIGN.CENTER, 'Consolas')
    add_text_box(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4),
                 'Configurable verification framework for extreme weather scorecard generation',
                 15, False, '#4E5D78', PP_ALIGN.CENTER)

    # Separator line
    line = slide.shapes.add_connector(1,
        Inches(1.5), Inches(1.15), Inches(11.8), Inches(1.15))
    line.line.color.rgb = rgb('#D6DCE5')
    line.line.width = Pt(1.5)

    # ── INPUT CARDS (left) ───────────────────────────────────────
    info_box(Inches(0.3), Inches(1.4), Inches(2.8), Inches(1.6),
             '\u25B6  Forecasts',
             ['Model 1 vs Model 2', 'e.g. IFS oper vs AIFS', 'Det or Ens (50 mbr)'],
             '#3D5A80', '#E8EEF5', '#6BA3CC')

    info_box(Inches(0.3), Inches(3.2), Inches(2.8), Inches(1.6),
             '\u25CF  Observations',
             ['SYNOP stations', 'HDOBS (high-density)', 'via STVL / local .gpt'],
             '#EE6C4D', '#FFF0EC', '#E8947E')

    # ── CONFIG (top-right) ───────────────────────────────────────
    info_box(Inches(10.2), Inches(1.4), Inches(2.8), Inches(1.6),
             '\u2699  config.yaml',
             ['variable \u00B7 date range', 'thresholds \u00B7 scores', 'filters \u00B7 output'],
             '#7B2D8E', '#F3E8F9', '#C9A0DC')

    # ── BACKENDS (right) ─────────────────────────────────────────
    add_rounded_box(Inches(10.2), Inches(3.2), Inches(1.3), Inches(0.75),
                    '#E8F5E9', '#81C784')
    add_text_box(Inches(10.2), Inches(3.25), Inches(1.3), Inches(0.35),
                 'Local', 12, True, '#2D6A4F', PP_ALIGN.CENTER)
    add_text_box(Inches(10.2), Inches(3.55), Inches(1.3), Inches(0.25),
                 'GRIB on disk', 9, False, '#2D6A4F', PP_ALIGN.CENTER)

    add_rounded_box(Inches(11.7), Inches(3.2), Inches(1.3), Inches(0.75),
                    '#F3E5F5', '#CE93D8')
    add_text_box(Inches(11.7), Inches(3.25), Inches(1.3), Inches(0.35),
                 'Quaver', 12, True, '#6A1B9A', PP_ALIGN.CENTER)
    add_text_box(Inches(11.7), Inches(3.55), Inches(1.3), Inches(0.25),
                 'MARS / VTB', 9, False, '#6A1B9A', PP_ALIGN.CENTER)

    add_text_box(Inches(11.35), Inches(3.35), Inches(0.5), Inches(0.3),
                 '/', 16, True, '#4E5D78', PP_ALIGN.CENTER)

    # ── MAIN PIPELINE ────────────────────────────────────────────
    px = Inches(4.0)
    pw = Inches(3.2)
    ph = Inches(0.7)
    y = Inches(1.5)
    g = Inches(0.18)

    step_box(px, y, 1, 'Read Data', 'Load forecast & obs fields', pw, ph)

    # Arrows from inputs
    add_arrow(Inches(3.1), Inches(2.1), px, Inches(1.85), '#3D5A80')
    add_arrow(Inches(3.1), Inches(3.9), px, Inches(1.85), '#EE6C4D')
    # Arrow from config
    add_arrow(Inches(10.2), Inches(2.1), px + pw, Inches(1.85), '#7B2D8E')

    y += ph + g
    step_box(px, y, 2, 'Pre-process', 'Lapse-rate, wind, units', pw, ph)
    add_arrow(px + pw//2, y - g, px + pw//2, y)

    y += ph + g
    step_box(px, y, 3, 'Extract Points', 'Nearest gridpoint / aligned()', pw, ph)
    add_arrow(px + pw//2, y - g, px + pw//2, y)
    # Output label
    add_text_box(px + pw + Inches(0.15), y + Inches(0.15),
                 Inches(2.5), Inches(0.35),
                 '\u2192 Parquet files (per day)', 11, True, '#2D6A4F')

    y += ph + g
    step_box(px, y, 4, 'Filter & QC', 'Season, orography, outliers', pw, ph)
    add_arrow(px + pw//2, y - g, px + pw//2, y)

    y += ph + g
    step_box(px, y, 5, 'Threshold', 'Fixed / dataset / station clim.', pw, ph)
    add_arrow(px + pw//2, y - g, px + pw//2, y)

    # ── BRANCH PANELS ────────────────────────────────────────────
    det_l = Inches(0.5)
    det_w = Inches(4.2)
    ens_l = Inches(5.5)
    ens_w = Inches(7.5)
    br_t = Inches(5.1)
    br_h = Inches(2.05)

    # Det panel
    det_panel = add_rounded_box(det_l, br_t, det_w, br_h, '#DBEAF5', '#6BA3CC', Pt(2))
    add_text_box(det_l + Inches(0.2), br_t + Inches(0.05),
                 det_w - Inches(0.4), Inches(0.35),
                 'DETERMINISTIC', 16, True, '#0B3954', PP_ALIGN.CENTER)

    branch_box(det_l + Inches(0.2), br_t + Inches(0.48),
               '\u2465 Scores', 'ETS \u00B7 PSS \u00B7 POD \u00B7 FAR \u00B7 twMAE \u00B7 bias',
               det_w - Inches(0.4), '#0B3954', '#6BA3CC')

    branch_box(det_l + Inches(0.2), br_t + Inches(1.25),
               '\u2466 Bootstrap', 'Paired significance testing',
               det_w - Inches(0.4), '#0B3954', '#6BA3CC')

    # Ens panel
    add_rounded_box(ens_l, br_t, ens_w, br_h, '#FDF0E2', '#D4A574', Pt(2))
    add_text_box(ens_l + Inches(0.2), br_t + Inches(0.05),
                 ens_w - Inches(0.4), Inches(0.35),
                 'ENSEMBLE', 16, True, '#7A4A1E', PP_ALIGN.CENTER)

    branch_box(ens_l + Inches(0.2), br_t + Inches(0.48),
               '\u2465 Scores', 'twCRPS \u00B7 Brier \u00B7 BSS \u00B7 QS \u00B7 twMAE \u00B7 tw SSR',
               ens_w - Inches(0.4), '#6D4C28', '#D4A574')

    branch_box(ens_l + Inches(0.2), br_t + Inches(1.25),
               '\u2466 Bootstrap', 'Paired significance testing',
               ens_w - Inches(0.4), '#6D4C28', '#D4A574')

    # Branch arrows from step 5
    step5_bottom = y + ph
    add_arrow(px + pw//4, step5_bottom, det_l + det_w//2, br_t, '#6BA3CC')
    add_arrow(px + pw*3//4, step5_bottom, ens_l + ens_w//2, br_t, '#D4A574')

    # ── SAVE & PLOT ──────────────────────────────────────────────
    bot_y = Inches(7.5) - Inches(0.95)
    step_box(Inches(3.0), bot_y, 8, 'Save Results', 'CSV per season \u00D7 orog', Inches(3.0), Inches(0.65))
    step_box(Inches(6.8), bot_y, 9, 'Plot Scorecards', 'Heatmaps \u00B7 panels', Inches(3.2), Inches(0.65))

    # Arrows from branches to save
    add_arrow(det_l + det_w//2, br_t + br_h, Inches(4.5), bot_y, '#6BA3CC')
    add_arrow(ens_l + ens_w//2, br_t + br_h, Inches(4.5), bot_y, '#D4A574')
    # Save → Plot
    add_arrow(Inches(6.0), bot_y + Inches(0.32), Inches(6.8), bot_y + Inches(0.32))

    # Output box
    add_rounded_box(Inches(10.3), bot_y, Inches(2.7), Inches(0.65),
                    '#D8F3DC', '#52B788')
    add_text_box(Inches(10.3), bot_y + Inches(0.02),
                 Inches(2.7), Inches(0.30),
                 '\u2192 Scorecard Heatmaps', 12, True, '#2D6A4F', PP_ALIGN.CENTER)
    add_text_box(Inches(10.3), bot_y + Inches(0.32),
                 Inches(2.7), Inches(0.25),
                 'PNG / PDF per score', 9.5, False, '#2D6A4F', PP_ALIGN.CENTER)
    add_arrow(Inches(10.0), bot_y + Inches(0.32), Inches(10.3), bot_y + Inches(0.32), '#52B788')

    # ── VARIABLE BADGES (bottom-left) ────────────────────────────
    vx = Inches(0.3)
    add_rounded_box(vx, bot_y, Inches(2.3), Inches(0.65), '#FFFFFF', '#CDD5E0')
    add_text_box(vx + Inches(0.1), bot_y + Inches(0.02),
                 Inches(2.1), Inches(0.25),
                 'Variables:', 11, True, '#1A1F36', PP_ALIGN.LEFT)
    for i, v in enumerate(['2t', '10ff', 'tp24']):
        bx = vx + Inches(0.1) + Inches(i * 0.72)
        shp = add_rounded_box(bx, bot_y + Inches(0.30), Inches(0.62), Inches(0.28),
                               '#1B4965', None)
        add_text_box(bx, bot_y + Inches(0.30), Inches(0.62), Inches(0.28),
                     v, 10, True, '#FFFFFF', PP_ALIGN.CENTER, 'Consolas')

    out = './plots/scorecards4extremes_overview.pptx'
    prs.save(out)
    print(f'PPTX saved: {out}')


if __name__ == '__main__':
    make_figure()
    make_pptx()
